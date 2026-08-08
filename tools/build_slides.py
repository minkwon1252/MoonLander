#!/usr/bin/env python3
"""Builds the Tech Race pre-game quick guide deck (PPTX).

Bright conference style, very large type, minimal text, Selene logo on every slide.
Regenerate with:  python3 tools/build_slides.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
LOGO = "/home/min/SNU/STEM/3rd_project_Selene_Program/selene_logo.png"
OUT = os.path.join(ROOT, "Tech_Race_Quick_Guide.pptx")

# --- palette: bright, high contrast, projector friendly ---
BG        = RGBColor(0xFA, 0xFB, 0xFF)   # near-white page
INK       = RGBColor(0x10, 0x18, 0x33)   # near-black navy text
NAVY      = RGBColor(0x1B, 0x2A, 0x5B)   # headings
GOLD      = RGBColor(0xD9, 0x9A, 0x1E)   # accent rule / highlights
SLATE     = RGBColor(0x53, 0x5D, 0x7E)   # secondary text
CARD      = RGBColor(0xEE, 0xF1, 0xFA)   # tint block
CARD_EDGE = RGBColor(0xD3, 0xDA, 0xEE)

W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_bg(slide, color=BG):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    solid(s, color)
    s.shadow.inherit = False
    # push to back
    slide.shapes._spTree.remove(s._element)
    slide.shapes._spTree.insert(2, s._element)
    return s


def add_logo(slide, big=False):
    """Selene logo — large and centered on the title slide, small top-right elsewhere."""
    if big:
        h = Inches(2.0)
        slide.shapes.add_picture(LOGO, int((W - h) / 2), Inches(0.85), height=h)
    else:
        h = Inches(0.72)
        slide.shapes.add_picture(LOGO, int(W - h - Inches(0.5)), Inches(0.38), height=h)


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, size, color=INK, bold=False, space_after=10, align=PP_ALIGN.LEFT,
         first=False, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return p


def heading(slide, title, kicker=None):
    """Standard slide header: small gold kicker, big navy title, gold rule."""
    from pptx.enum.shapes import MSO_SHAPE
    y = Inches(0.45)
    if kicker:
        tf = textbox(slide, Inches(0.75), y, Inches(9.5), Inches(0.4))
        para(tf, kicker.upper(), 15, GOLD, bold=True, space_after=0, first=True)
        y = Inches(0.82)
    tf = textbox(slide, Inches(0.75), y, Inches(10.2), Inches(1.0))
    para(tf, title, 46, NAVY, bold=True, space_after=0, first=True)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.82),
                                  Inches(1.5), Pt(5))
    solid(rule, GOLD)
    rule.shadow.inherit = False


def bullets(slide, items, size=30, top=2.35, gap=14, color=INK, left=0.95, width=11.4):
    tf = textbox(slide, Inches(left), Inches(top), Inches(width), Inches(4.4))
    for i, it in enumerate(items):
        para(tf, it, size, color, space_after=gap, first=(i == 0))


def new(kicker=None, title=None, big_logo=False):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_logo(s, big=big_logo)
    if title:
        heading(s, title, kicker)
    return s


# ------------------------------------------------------------------ 1. Title
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_logo(s, big=True)
tf = textbox(s, Inches(0.8), Inches(3.05), Inches(11.7), Inches(1.4))
para(tf, "TECH RACE", 78, NAVY, bold=True, space_after=4, align=PP_ALIGN.CENTER, first=True)
from pptx.enum.shapes import MSO_SHAPE
rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int((W - Inches(2.2)) / 2), Inches(4.52),
                          Inches(2.2), Pt(5))
solid(rule, GOLD)
rule.shadow.inherit = False
tf = textbox(s, Inches(0.8), Inches(4.85), Inches(11.7), Inches(1.2))
para(tf, "International Cooperation & Competition", 30, INK, space_after=14,
     align=PP_ALIGN.CENTER, first=True)
para(tf, "Selene Program", 22, SLATE, align=PP_ALIGN.CENTER)

# ------------------------------------------------------------------ 2. Why We Play
s = new("01", "Why We Play")
bullets(s, [
    "Technology is never only technical.",
    "Every decision is political, economic and social.",
    "Here, you defend your reasoning out loud.",
], size=32, top=2.6, gap=30)

# ------------------------------------------------------------------ 3. Your Role
s = new("02", "Your Role")
bullets(s, [
    "Your team leads a national innovation system.",
    "Choose one technology goal.",
    "Develop it faster — without breaking your country.",
], size=32, top=2.6, gap=30)

# ------------------------------------------------------------------ 4. The 8 Balances
s = new("03", "The 8 Balances")
stats = [
    "Treasury", "Energy / Compute",
    "Political Support", "Public Welfare",
    "R&D Capacity", "International Reputation",
    "Security / Sovereignty", "Environment",
]
CW, CH = Inches(5.7), Inches(0.95)
for i, name in enumerate(stats):
    col, row = i % 2, i // 2
    x = Inches(0.83) + col * (CW + Inches(0.3))
    y = Inches(2.25) + row * (CH + Inches(0.2))
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, CW, CH)
    solid(box, CARD)
    box.line.color.rgb = CARD_EDGE
    box.line.width = Pt(1)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = name
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = NAVY; r.font.name = "Calibri"

tf = textbox(s, Inches(0.85), Inches(6.85), Inches(11.4), Inches(0.5))
para(tf, "Let any of these collapse and your technology stops advancing.", 20, GOLD,
     bold=True, align=PP_ALIGN.CENTER, first=True)

# ------------------------------------------------------------------ 5. How a Round Works
s = new("04", "How a Round Works")
steps = [
    ("1", "A global event appears"),
    ("2", "Each team gets a policy card"),
    ("3", "Your team explains its reasoning"),
    ("4", "Swipe left or right"),
    ("5", "Stats and roadmap update"),
]
BW = Inches(2.24)
for i, (n, label) in enumerate(steps):
    x = Inches(0.72) + i * (BW + Inches(0.16))
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.75), BW, Inches(2.3))
    solid(box, CARD)
    box.line.color.rgb = CARD_EDGE
    box.line.width = Pt(1)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.18)
    # anchor to top so the big step numbers line up across cards regardless of label length
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(4)
    r = p.add_run(); r.text = n
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = GOLD; r.font.name = "Calibri"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = label
    r2.font.size = Pt(19); r2.font.color.rgb = INK; r2.font.name = "Calibri"

tf = textbox(s, Inches(0.75), Inches(5.5), Inches(11.6), Inches(0.6))
para(tf, "8 rounds. About two hours.", 26, NAVY, bold=True, align=PP_ALIGN.CENTER, first=True)

# ------------------------------------------------------------------ 6. Global Events
s = new("05", "Global Events")
bullets(s, [
    "Every round has one random event.",
    "Crisis, opportunity, resource conflict or pressure.",
], size=30, top=2.5, gap=24)
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(4.75),
                         Inches(11.5), Inches(1.7))
solid(box, CARD)
box.line.color.rgb = GOLD
box.line.width = Pt(2.5)
box.shadow.inherit = False
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_right = Inches(0.4)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Read the event guide before deciding."
r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = NAVY; r.font.name = "Calibri"
p2 = tf.add_paragraph()
r2 = p2.add_run()
r2.text = '"Energy Shock — energy-heavy strategies become riskier this round."'
r2.font.size = Pt(21); r2.font.italic = True; r2.font.color.rgb = SLATE; r2.font.name = "Calibri"

# ------------------------------------------------------------------ 7. Development Roadmap
s = new("06", "Development Roadmap")
bullets(s, [
    "Your technology starts at 0%.",
    "Progress comes step by step, round by round.",
], size=32, top=2.4, gap=22)
# progress strip
seg_w = Inches(2.05)
for i in range(5):
    x = Inches(0.95) + i * (seg_w + Inches(0.22))
    filled = i < 2
    seg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.35), seg_w, Inches(0.72))
    solid(seg, GOLD if filled else CARD)
    seg.line.color.rgb = GOLD if filled else CARD_EDGE
    seg.line.width = Pt(1.5)
    seg.shadow.inherit = False
    tf = seg.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = f"Stage {i+1}"
    r.font.size = Pt(17); r.font.bold = True
    r.font.color.rgb = RGBColor(0x2A, 0x1D, 0x00) if filled else SLATE
    r.font.name = "Calibri"

tf = textbox(s, Inches(0.95), Inches(5.5), Inches(11.4), Inches(1.0))
para(tf, "If your national balance collapses,", 28, INK, space_after=2,
     align=PP_ALIGN.CENTER, first=True)
para(tf, "progress stops — or goes backwards.", 30, RGBColor(0xB4, 0x3A, 0x2A), bold=True,
     align=PP_ALIGN.CENTER)

# ------------------------------------------------------------------ 8. Winning
s = new("07", "Winning")
bullets(s, [
    "Fast technology growth matters.",
    "But an unstable country cannot sustain innovation.",
], size=30, top=2.4, gap=20)
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(4.75),
                         Inches(11.5), Inches(1.5))
solid(box, NAVY)
box.line.fill.background()
box.shadow.inherit = False
tf = box.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Best team  =  strong technology  +  balanced society  +  strategic judgement"
r.font.size = Pt(22); r.font.bold = True
r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); r.font.name = "Calibri"

# ------------------------------------------------------------------ 9. Final Message
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_logo(s)
tf = textbox(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(4.6))
para(tf, "Think like engineers.", 46, NAVY, bold=True, space_after=18,
     align=PP_ALIGN.CENTER, first=True)
para(tf, "Decide like policymakers.", 46, NAVY, bold=True, space_after=18, align=PP_ALIGN.CENTER)
para(tf, "Negotiate like diplomats.", 46, NAVY, bold=True, space_after=30, align=PP_ALIGN.CENTER)
para(tf, "Explain your choices.", 40, GOLD, bold=True, align=PP_ALIGN.CENTER)
rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int((W - Inches(2.6)) / 2), Inches(5.82),
                          Inches(2.6), Pt(5))
solid(rule, GOLD)
rule.shadow.inherit = False
tf = textbox(s, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.5))
para(tf, "Tech Race  ·  Selene Program", 18, SLATE, align=PP_ALIGN.CENTER, first=True)

prs.save(OUT)
print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
